#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
血压守护 - PPT汇报材料生成脚本
10页以内，重点介绍功能点、社会必要性、社会价值
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

class BPGuardianPPT:
    """血压守护PPT生成器"""
    
    def __init__(self):
        self.prs = Presentation()
        # 设置16:9宽屏
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # 主题色
        self.COLOR_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)    # 深蓝
        self.COLOR_SECONDARY = RGBColor(0x2E, 0x75, 0xB6)  # 中蓝
        self.COLOR_ACCENT = RGBColor(0x44, 0x72, 0xC4)     # 亮蓝
        self.COLOR_GREEN = RGBColor(0x4C, 0xAF, 0x50)      # 绿色
        self.COLOR_ORANGE = RGBColor(0xED, 0x7D, 0x31)     # 橙色
        self.COLOR_RED = RGBColor(0xC0, 0x00, 0x00)        # 红色
        self.COLOR_LIGHT_BG = RGBColor(0xF2, 0xF7, 0xFB)   # 浅蓝背景
        self.COLOR_DARK_TEXT = RGBColor(0x33, 0x33, 0x33)  # 深灰文字
        self.COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)      # 白色
    
    def add_background(self, slide, color):
        """添加背景色"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def add_shape(self, slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
        """添加形状"""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    def add_text_box(self, slide, left, top, width, height, text, font_size=18, 
                     bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or self.COLOR_DARK_TEXT
        p.font.name = font_name
        p.alignment = alignment
        return txBox
    
    def add_bullet_list(self, slide, left, top, width, height, items, font_size=16, 
                       color=None, bullet_char='✓'):
        """添加项目列表"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"  {bullet_char}  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color or self.COLOR_DARK_TEXT
            p.font.name = '微软雅黑'
            p.space_after = Pt(8)
        
        return txBox
    
    def add_info_card(self, slide, left, top, width, height, title, content, 
                     bg_color=None, title_color=None):
        """添加信息卡片"""
        # 卡片背景
        card = self.add_shape(slide, left, top, width, height, bg_color or self.COLOR_LIGHT_BG)
        
        # 标题
        self.add_text_box(slide, left + Inches(0.3), top + Inches(0.15), 
                         width - Inches(0.6), Inches(0.5),
                         title, font_size=16, bold=True, 
                         color=title_color or self.COLOR_PRIMARY)
        
        # 内容
        self.add_text_box(slide, left + Inches(0.3), top + Inches(0.6), 
                         width - Inches(0.6), height - Inches(0.7),
                         content, font_size=14, color=self.COLOR_DARK_TEXT)
        
        return card
    
    # ==================== 第1页：封面 ====================
    def slide_1_cover(self):
        """封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白版式
        self.add_background(slide, self.COLOR_PRIMARY)
        
        # 顶部装饰条
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(0.15), self.COLOR_ACCENT)
        
        # 底部装饰条
        self.add_shape(slide, Inches(0), Inches(7.35), 
                      self.prs.slide_width, Inches(0.15), self.COLOR_ACCENT)
        
        # 主标题
        self.add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.5),
                         '血压守护', font_size=54, bold=True, 
                         color=self.COLOR_WHITE, alignment=PP_ALIGN.CENTER)
        
        # 副标题
        self.add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(1),
                         '智能服药提醒与血压记录管理系统', 
                         font_size=28, color=self.COLOR_WHITE, alignment=PP_ALIGN.CENTER)
        
        # 分隔线
        self.add_shape(slide, Inches(5), Inches(4.5), Inches(3.3), Inches(0.05), 
                      self.COLOR_ACCENT)
        
        # 汇报信息
        self.add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
                         '项目汇报材料', font_size=20, 
                         color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)
        
        self.add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
                         '让每一位高血压患者都能得到及时、准确、专业的健康管理服务',
                         font_size=16, color=RGBColor(0x90, 0xCA, 0xF9), 
                         alignment=PP_ALIGN.CENTER)
        
        # 日期
        self.add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5),
                         '2026年5月', font_size=14, 
                         color=RGBColor(0x64, 0xB5, 0xF6), alignment=PP_ALIGN.CENTER)
    
    # ==================== 第2页：社会背景 ====================
    def slide_2_background(self):
        """社会背景与必要性"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLOR_WHITE)
        
        # 顶部标题栏
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(1.2), self.COLOR_PRIMARY)
        
        self.add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
                         '一、社会背景与健康需求', font_size=32, bold=True, 
                         color=self.COLOR_WHITE)
        
        # 左侧：数据卡片
        data_cards = [
            ('2.7亿+', '中国高血压患者', '18岁以上患病率23.2%'),
            ('15.3%', '血压控制率', '远低于发达国家水平'),
            ('30-40%', '服药依从性', '超60%患者不能按时服药'),
        ]
        
        for i, (number, title, desc) in enumerate(data_cards):
            left = Inches(0.8) + Inches(3.8) * i
            top = Inches(1.8)
            
            # 卡片
            card = self.add_shape(slide, left, top, Inches(3.5), Inches(2), self.COLOR_LIGHT_BG)
            
            # 数字
            self.add_text_box(slide, left, top + Inches(0.2), Inches(3.5), Inches(0.8),
                             number, font_size=40, bold=True, 
                             color=self.COLOR_RED, alignment=PP_ALIGN.CENTER)
            
            # 标题
            self.add_text_box(slide, left, top + Inches(1), Inches(3.5), Inches(0.5),
                             title, font_size=18, bold=True, 
                             color=self.COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)
            
            # 描述
            self.add_text_box(slide, left, top + Inches(1.4), Inches(3.5), Inches(0.5),
                             desc, font_size=12, 
                             color=self.COLOR_DARK_TEXT, alignment=PP_ALIGN.CENTER)
        
        # 右侧：痛点分析
        self.add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.6),
                         '核心痛点', font_size=24, bold=True, 
                         color=self.COLOR_PRIMARY)
        
        pain_points = [
            '忘记服药 — 手机闹钟无法追踪是否真正服药',
            '记录混乱 — 纸质记录容易丢失，无法分析趋势',
            '预警缺失 — 血压异常时无法及时发现和处理',
            '家属担忧 — 异地子女无法实时了解父母健康',
            '就医低效 — 缺乏完整数据，医生难以准确判断'
        ]
        
        self.add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.5),
                            pain_points, font_size=16, bullet_char='●')
        
        # 右侧：必要性说明
        self.add_shape(slide, Inches(7.5), Inches(4.2), Inches(5.3), Inches(2.8), self.COLOR_LIGHT_BG)
        
        self.add_text_box(slide, Inches(7.8), Inches(4.3), Inches(4.8), Inches(0.5),
                         '为什么需要"血压守护"？', font_size=20, bold=True, 
                         color=self.COLOR_PRIMARY)
        
        necessity = [
            '国家政策大力支持"互联网+医疗健康"',
            '老龄化加速，慢性病管理需求激增',
            '传统健康管理方式已无法满足需求',
            '数字化工具可提升服药依从性2-3倍',
            '降低高血压并发症，减少医疗支出'
        ]
        
        self.add_bullet_list(slide, Inches(7.8), Inches(4.9), Inches(4.8), Inches(2),
                            necessity, font_size=15, color=self.COLOR_DARK_TEXT,
                            bullet_char='★')
    
    # ==================== 第3页：解决方案 ====================
    def slide_3_solution(self):
        """解决方案概述"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLOR_WHITE)
        
        # 顶部标题栏
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(1.2), self.COLOR_PRIMARY)
        
        self.add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
                         '二、解决方案 — 血压守护', font_size=32, bold=True, 
                         color=self.COLOR_WHITE)
        
        # 产品定位
        self.add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
                         '产品定位：专注高血压患者及其家庭成员的智能健康管理平台',
                         font_size=20, bold=True, color=self.COLOR_PRIMARY,
                         alignment=PP_ALIGN.CENTER)
        
        # 四大核心能力卡片
        cards = [
            ('💊', '智能服药管理', '多药物提醒\n服药确认追踪\n漏服预警\n库存管理'),
            ('📊', '血压数据管理', '快速记录血压\n趋势图表分析\n健康评估\nPDF报告导出'),
            ('🔔', '分级预警系统', '漏服实时预警\n血压异常预警\n三级预警机制\n多渠道通知'),
            ('👨\u200D👩\u200D👧', '家庭协作管理', '多成员管理\n数据共享查看\n代录数据\n协作日志'),
        ]
        
        for i, (icon, title, content) in enumerate(cards):
            left = Inches(0.8) + Inches(3.1) * i
            top = Inches(2.5)
            
            # 卡片背景
            self.add_shape(slide, left, top, Inches(2.9), Inches(4.5), self.COLOR_LIGHT_BG)
            
            # 图标
            self.add_text_box(slide, left, top + Inches(0.3), Inches(2.9), Inches(0.8),
                             icon, font_size=48, alignment=PP_ALIGN.CENTER)
            
            # 标题
            self.add_text_box(slide, left, top + Inches(1.2), Inches(2.9), Inches(0.6),
                             title, font_size=18, bold=True, 
                             color=self.COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)
            
            # 内容
            txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(2), 
                                            Inches(2.5), Inches(2.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            lines = content.split('\n')
            for j, line in enumerate(lines):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f'  ✓  {line}'
                p.font.size = Pt(14)
                p.font.color.rgb = self.COLOR_DARK_TEXT
                p.font.name = '微软雅黑'
                p.space_after = Pt(6)
        
        # 底部说明
        self.add_shape(slide, Inches(0), Inches(7.2), 
                      self.prs.slide_width, Inches(0.3), self.COLOR_PRIMARY)
    
    # ==================== 第4页：核心功能-服药管理 ====================
    def slide_4_medication(self):
        """服药管理功能详情"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLOR_WHITE)
        
        # 顶部标题栏
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(1.2), self.COLOR_SECONDARY)
        
        self.add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
                         '三、核心功能 — 智能服药管理', font_size=32, bold=True, 
                         color=self.COLOR_WHITE)
        
        # 左侧功能列表
        features = [
            '多药物管理：支持添加50种药物，分别设置服用时间',
            '智能提醒：APP推送、短信、电话多种提醒方式',
            '服药确认：一键确认，记录准确服药时间',
            '漏服预警：超时未服自动预警并通知家属',
            '库存管理：自动计算剩余药量，提前提醒购药',
            '服药历史：完整记录，支持追溯和统计',
        ]
        
        self.add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(3.5),
                            features, font_size=18, bullet_char='✓',
                            color=self.COLOR_DARK_TEXT)
        
        # 右侧：效果数据
        self.add_shape(slide, Inches(7.5), Inches(1.6), Inches(5.3), Inches(3.5), self.COLOR_LIGHT_BG)
        
        self.add_text_box(slide, Inches(7.8), Inches(1.7), Inches(4.8), Inches(0.5),
                         '应用效果', font_size=22, bold=True, color=self.COLOR_PRIMARY)
        
        effects = [
            '服药依从性：30% → 80-90%（提升2-3倍）',
            '漏服率：70% → <10%',
            '预警响应时间：<5分钟',
            '用户满意度：NPS评分 ≥ 50',
            '日均服药记录：3次/用户',
        ]
        
        self.add_bullet_list(slide, Inches(7.8), Inches(2.3), Inches(4.8), Inches(2.5),
                            effects, font_size=16, bullet_char='📈',
                            color=self.COLOR_DARK_TEXT)
        
        # 底部：业务流程
        self.add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
                         '业务流程', font_size=20, bold=True, color=self.COLOR_PRIMARY)
        
        # 流程图
        steps = ['添加药物', '设置时间', '接收提醒', '确认服药', '记录追踪']
        for i, step in enumerate(steps):
            left = Inches(1) + Inches(2.4) * i
            # 步骤框
            self.add_shape(slide, left, Inches(5.9), Inches(1.8), Inches(0.7), self.COLOR_ACCENT)
            self.add_text_box(slide, left, Inches(5.95), Inches(1.8), Inches(0.6),
                             step, font_size=16, bold=True, 
                             color=self.COLOR_WHITE, alignment=PP_ALIGN.CENTER)
            
            # 箭头
            if i < len(steps) - 1:
                self.add_text_box(slide, left + Inches(1.85), Inches(5.95), Inches(0.5), Inches(0.6),
                                 '→', font_size=20, bold=True, 
                                 color=self.COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)
    
    # ==================== 第5页：核心功能-血压管理 ====================
    def slide_5_bp_management(self):
        """血压管理功能详情"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLOR_WHITE)
        
        # 顶部标题栏
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(1.2), self.COLOR_SECONDARY)
        
        self.add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
                         '四、核心功能 — 血压数据管理', font_size=32, bold=True, 
                         color=self.COLOR_WHITE)
        
        # 左侧功能
        features = [
            '快捷录入：三步完成血压记录（收缩压/舒张压/心率）',
            '趋势分析：可视化折线图展示7天/30天/90天趋势',
            '数据统计：自动计算平均值、最高值、最低值',
            '健康评估：基于近期数据评估健康状况等级',
            '异常标注：超出正常范围的数据高亮显示',
            '报告导出：生成PDF报告，方便就医使用',
        ]
        
        self.add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(3.5),
                            features, font_size=18, bullet_char='✓',
                            color=self.COLOR_DARK_TEXT)
        
        # 右侧：血压标准
        self.add_shape(slide, Inches(7.5), Inches(1.6), Inches(5.3), Inches(3.5), self.COLOR_LIGHT_BG)
        
        self.add_text_box(slide, Inches(7.8), Inches(1.7), Inches(4.8), Inches(0.5),
                         '血压参考标准', font_size=22, bold=True, color=self.COLOR_PRIMARY)
        
        # 标准表格数据
        standards = [
            ['指标', '正常范围', '偏高', '危险'],
            ['收缩压', '90-140 mmHg', '140-160', '≥160'],
            ['舒张压', '60-90 mmHg', '90-110', '≥110'],
            ['心率', '60-100 次/分', '100-120', '≥120'],
        ]
        
        # 手动创建表格
        table_left = Inches(7.8)
        table_top = Inches(2.4)
        col_widths = [Inches(1.3), Inches(1.4), Inches(1.2), Inches(1.2)]
        
        for row_idx, row_data in enumerate(standards):
            left = table_left
            for col_idx, cell_data in enumerate(row_data):
                # 背景色
                if row_idx == 0:
                    bg_color = self.COLOR_PRIMARY
                    text_color = self.COLOR_WHITE
                elif col_idx == 3:
                    bg_color = RGBColor(0xFF, 0xE6, 0xE6)
                    text_color = self.COLOR_RED
                else:
                    bg_color = self.COLOR_WHITE
                    text_color = self.COLOR_DARK_TEXT
                
                # 单元格
                self.add_shape(slide, left, table_top + Inches(row_idx * 0.5), 
                              col_widths[col_idx], Inches(0.5), bg_color)
                
                self.add_text_box(slide, left, table_top + Inches(row_idx * 0.5), 
                                 col_widths[col_idx], Inches(0.5),
                                 cell_data, font_size=14, bold=(row_idx==0),
                                 color=text_color, alignment=PP_ALIGN.CENTER)
                
                left += col_widths[col_idx]
        
        # 底部说明
        self.add_text_box(slide, Inches(7.8), Inches(4.8), Inches(4.8), Inches(0.5),
                         '数据来源：《中国高血压防治指南（2023年版）》', 
                         font_size=12, color=RGBColor(0x99, 0x99, 0x99))
        
        # 底部：使用价值
        self.add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
                         '使用价值', font_size=20, bold=True, color=self.COLOR_PRIMARY)
        
        values = [
            '帮助用户了解血压变化规律，及时调整生活方式',
            '为医生诊断提供完整、准确的数据支撑',
            '通过趋势分析，提前发现健康风险',
            'PDF报告导出，就诊效率提升50%'
        ]
        
        self.add_bullet_list(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.5),
                            values, font_size=15, bullet_char='💡',
                            color=self.COLOR_DARK_TEXT)
    
    # ==================== 第6页：核心功能-预警与家庭 ====================
    def slide_6_alert_family(self):
        """预警系统与家庭协作"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLOR_WHITE)
        
        # 顶部标题栏
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(1.2), self.COLOR_SECONDARY)
        
        self.add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
                         '五、核心功能 — 预警系统 & 家庭协作', font_size=32, bold=True, 
                         color=self.COLOR_WHITE)
        
        # 左侧：预警系统
        self.add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
                         '🔔 分级预警系统', font_size=24, bold=True, 
                         color=self.COLOR_PRIMARY)
        
        # 预警级别
        alerts = [
            ('⚠️ 警告', '漏服超过2小时', RGBColor(0xFF, 0xF3, 0xE0)),
            ('🔶 紧急', '收缩压≥160或舒张压≥110', RGBColor(0xFF, 0xE0, 0xB2)),
            ('🔴 危急', '收缩压≥180或舒张压≥120', RGBColor(0xFF, 0xCD, 0xCD)),
        ]
        
        for i, (level, condition, bg_color) in enumerate(alerts):
            top = Inches(2.3) + Inches(1.2) * i
            
            self.add_shape(slide, Inches(0.8), top, Inches(5.5), Inches(1), bg_color)
            self.add_text_box(slide, Inches(1), top + Inches(0.1), Inches(2), Inches(0.4),
                             level, font_size=18, bold=True, 
                             color=self.COLOR_DARK_TEXT)
            self.add_text_box(slide, Inches(3), top + Inches(0.1), Inches(3), Inches(0.4),
                             f'触发条件：{condition}', font_size=14, 
                             color=self.COLOR_DARK_TEXT)
            self.add_text_box(slide, Inches(1), top + Inches(0.55), Inches(5), Inches(0.4),
                             '通知方式：APP推送 + 短信 + 家属通知', font_size=12, 
                             color=RGBColor(0x66, 0x66, 0x66))
        
        # 右侧：家庭协作
        self.add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.6),
                         '👨\u200D👩\u200D👧 家庭协作管理', font_size=24, bold=True, 
                         color=self.COLOR_PRIMARY)
        
        family_features = [
            '添加多个家庭成员，分别管理健康数据',
            '家属可实时查看患者血压和服药情况',
            '支持家属代为记录血压和确认服药',
            '异常情况自动通知所有家庭成员',
            '完整协作日志，操作可追溯',
            '解决异地照顾难题，增强家庭关爱',
        ]
        
        self.add_bullet_list(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(3),
                            family_features, font_size=16, bullet_char='✓',
                            color=self.COLOR_DARK_TEXT)
        
        # 底部：社会价值
        self.add_shape(slide, Inches(0), Inches(6.3), Inches(13.333), Inches(1.2), self.COLOR_LIGHT_BG)
        
        self.add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
                         '💝 社会价值', font_size=18, bold=True, color=self.COLOR_PRIMARY)
        
        self.add_text_box(slide, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.5),
                         '让子女无论身在何处，都能随时了解父母健康状况，减少担忧，增强家庭关爱',
                         font_size=14, color=self.COLOR_DARK_TEXT, alignment=PP_ALIGN.CENTER)
    
    # ==================== 第7页：技术架构 ====================
    def slide_7_tech(self):
        """技术架构"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLOR_WHITE)
        
        # 顶部标题栏
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(1.2), self.COLOR_PRIMARY)
        
        self.add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
                         '六、技术架构与安全保障', font_size=32, bold=True, 
                         color=self.COLOR_WHITE)
        
        # 技术架构四层
        layers = [
            ('客户端层', '手机浏览器 / 微信内嵌 / PWA应用 / 小程序', self.COLOR_ACCENT),
            ('接入层', 'Nginx + HTTPS / 负载均衡 / 静态资源', self.COLOR_SECONDARY),
            ('应用层', 'Flask后端 / Gunicorn / API路由 / 业务逻辑', self.COLOR_PRIMARY),
            ('数据层', 'MySQL数据库 / Redis缓存 / 文件存储 / 日志系统', RGBColor(0x15, 0x3E, 0x60)),
        ]
        
        for i, (layer_name, desc, color) in enumerate(layers):
            top = Inches(1.6) + Inches(1.2) * i
            
            # 层名称
            self.add_shape(slide, Inches(1.5), top, Inches(2.5), Inches(0.8), color)
            self.add_text_box(slide, Inches(1.5), top + Inches(0.15), Inches(2.5), Inches(0.5),
                             layer_name, font_size=18, bold=True, 
                             color=self.COLOR_WHITE, alignment=PP_ALIGN.CENTER)
            
            # 描述
            self.add_shape(slide, Inches(4.2), top, Inches(7.5), Inches(0.8), self.COLOR_LIGHT_BG)
            self.add_text_box(slide, Inches(4.4), top + Inches(0.15), Inches(7.1), Inches(0.5),
                             desc, font_size=16, color=self.COLOR_DARK_TEXT)
            
            # 箭头
            if i < len(layers) - 1:
                self.add_text_box(slide, Inches(6.3), top + Inches(0.85), Inches(0.8), Inches(0.35),
                                 '↓', font_size=24, bold=True, 
                                 color=color, alignment=PP_ALIGN.CENTER)
        
        # 右侧安全特性
        self.add_text_box(slide, Inches(0.8), Inches(5.6), Inches(5), Inches(0.5),
                         '🔒 安全保障', font_size=20, bold=True, 
                         color=self.COLOR_PRIMARY)
        
        security = [
            '传输加密：HTTPS/TLS 1.2+',
            '密码加密：PBKDF2 + Salt',
            'XSS防护：输入过滤 + 输出编码',
            'CSRF防护：Token验证机制',
            '速率限制：防止恶意请求'
        ]
        
        self.add_bullet_list(slide, Inches(0.8), Inches(6.1), Inches(5.5), Inches(1.2),
                            security, font_size=14, bullet_char='✓',
                            color=self.COLOR_DARK_TEXT)
        
        # 右侧性能指标
        self.add_text_box(slide, Inches(7), Inches(5.6), Inches(5.5), Inches(0.5),
                         '⚡ 性能指标', font_size=20, bold=True, 
                         color=self.COLOR_PRIMARY)
        
        perf = [
            '页面加载：≤2秒',
            'API响应：≤500ms',
            '并发支持：≥1000用户',
            '系统可用率：≥99.9%',
            '数据备份：每天，保留30天'
        ]
        
        self.add_bullet_list(slide, Inches(7), Inches(6.1), Inches(5.5), Inches(1.2),
                            perf, font_size=14, bullet_char='✓',
                            color=self.COLOR_DARK_TEXT)
    
    # ==================== 第8页：社会价值 ====================
    def slide_8_social_value(self):
        """社会价值"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLOR_WHITE)
        
        # 顶部标题栏
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(1.2), self.COLOR_GREEN)
        
        self.add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
                         '七、社会价值与意义', font_size=32, bold=True, 
                         color=self.COLOR_WHITE)
        
        # 四大社会价值
        values = [
            ('🏥', '提升健康管理水平', '覆盖2.7亿高血压患者\n服药依从性提升2-3倍\n血压控制率提升30%'),
            ('💰', '降低医疗费用', '减少并发症住院\n年节省医疗费¥5000-10000/人\n减轻医保基金压力'),
            ('👨\u200D👩\u200D👧', '增强家庭关爱', '解决异地照顾难题\n家属实时了解健康状况\n减少焦虑和担忧'),
            ('📊', '助力医疗体系', '提供完整院外数据\n提升医生诊断效率\n优化医疗资源配置'),
        ]
        
        for i, (icon, title, content) in enumerate(values):
            left = Inches(0.8) + Inches(3.1) * i
            top = Inches(1.8)
            
            # 卡片
            self.add_shape(slide, left, top, Inches(2.9), Inches(4), self.COLOR_LIGHT_BG)
            
            # 图标
            self.add_text_box(slide, left, top + Inches(0.3), Inches(2.9), Inches(0.8),
                             icon, font_size=48, alignment=PP_ALIGN.CENTER)
            
            # 标题
            self.add_text_box(slide, left, top + Inches(1.2), Inches(2.9), Inches(0.6),
                             title, font_size=16, bold=True, 
                             color=self.COLOR_GREEN, alignment=PP_ALIGN.CENTER)
            
            # 内容
            txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(2), 
                                            Inches(2.5), Inches(1.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            lines = content.split('\n')
            for j, line in enumerate(lines):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f'• {line}'
                p.font.size = Pt(13)
                p.font.color.rgb = self.COLOR_DARK_TEXT
                p.font.name = '微软雅黑'
                p.space_after = Pt(4)
        
        # 底部总结
        self.add_shape(slide, Inches(0), Inches(6.8), 
                      self.prs.slide_width, Inches(0.7), self.COLOR_GREEN)
        
        self.add_text_box(slide, Inches(1), Inches(6.85), Inches(11.333), Inches(0.6),
                         '💎 核心价值：让每一位高血压患者都能得到及时、准确、专业的健康管理服务',
                         font_size=20, bold=True, color=self.COLOR_WHITE, 
                         alignment=PP_ALIGN.CENTER)
    
    # ==================== 第9页：实施计划 ====================
    def slide_9_roadmap(self):
        """实施计划"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLOR_WHITE)
        
        # 顶部标题栏
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(1.2), self.COLOR_PRIMARY)
        
        self.add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8),
                         '八、实施计划与发展目标', font_size=32, bold=True, 
                         color=self.COLOR_WHITE)
        
        # 时间轴
        milestones = [
            ('2026-05', 'MVP完成', '核心功能开发完成\n功能/压力测试通过\nBUG修复完毕'),
            ('2026-06', '公测启动', '100名种子用户\n收集用户反馈\n优化迭代'),
            ('2026-07', '正式发布', '正式上线运营\n市场推广启动\n用户增长'),
            ('2026-12', '规模化', '用户破10万\n医生端开发\n企业版上线'),
        ]
        
        for i, (date, title, desc) in enumerate(milestones):
            left = Inches(1) + Inches(3.1) * i
            top = Inches(2)
            
            # 时间点
            self.add_shape(slide, left, top, Inches(2.8), Inches(0.6), self.COLOR_ACCENT)
            self.add_text_box(slide, left, top + Inches(0.1), Inches(2.8), Inches(0.4),
                             date, font_size=18, bold=True, 
                             color=self.COLOR_WHITE, alignment=PP_ALIGN.CENTER)
            
            # 连接线
            if i < len(milestones) - 1:
                self.add_shape(slide, left + Inches(2.8), top + Inches(0.2), 
                              Inches(0.3), Inches(0.2), self.COLOR_ACCENT)
            
            # 标题
            self.add_text_box(slide, left, top + Inches(0.8), Inches(2.8), Inches(0.5),
                             title, font_size=20, bold=True, 
                             color=self.COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)
            
            # 描述卡片
            self.add_shape(slide, left, top + Inches(1.4), Inches(2.8), Inches(2), self.COLOR_LIGHT_BG)
            
            txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.5), 
                                            Inches(2.4), Inches(1.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            lines = desc.split('\n')
            for j, line in enumerate(lines):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f'• {line}'
                p.font.size = Pt(13)
                p.font.color.rgb = self.COLOR_DARK_TEXT
                p.font.name = '微软雅黑'
                p.space_after = Pt(4)
        
        # 底部目标
        self.add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5),
                         '发展目标', font_size=22, bold=True, color=self.COLOR_PRIMARY)
        
        goals = [
            '2027年底：用户规模达到50万+',
            '付费转化率：5-8%',
            '年收入：¥500-800万',
            '日活跃用户：10万+',
        ]
        
        self.add_bullet_list(slide, Inches(0.8), Inches(6.3), Inches(5.5), Inches(1),
                            goals, font_size=16, bullet_char='🎯',
                            color=self.COLOR_DARK_TEXT)
        
        self.add_text_box(slide, Inches(7), Inches(6.3), Inches(5.5), Inches(1),
                         '让"血压守护"成为高血压管理领域的标杆产品，惠及千万家庭！',
                         font_size=16, bold=True, color=self.COLOR_ACCENT)
    
    # ==================== 第10页：结束页 ====================
    def slide_10_end(self):
        """结束页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLOR_PRIMARY)
        
        # 顶部装饰条
        self.add_shape(slide, Inches(0), Inches(0), 
                      self.prs.slide_width, Inches(0.15), self.COLOR_ACCENT)
        
        # 底部装饰条
        self.add_shape(slide, Inches(0), Inches(7.35), 
                      self.prs.slide_width, Inches(0.15), self.COLOR_ACCENT)
        
        # 感谢
        self.add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.5),
                         '感谢聆听', font_size=54, bold=True, 
                         color=self.COLOR_WHITE, alignment=PP_ALIGN.CENTER)
        
        # 分隔线
        self.add_shape(slide, Inches(5), Inches(3.2), Inches(3.3), Inches(0.05), 
                      self.COLOR_ACCENT)
        
        # 核心价值
        self.add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(1),
                         '让每一位高血压患者都能得到\n及时、准确、专业的健康管理服务',
                         font_size=24, color=RGBColor(0xBB, 0xDE, 0xFB), 
                         alignment=PP_ALIGN.CENTER)
        
        # 项目状态
        self.add_shape(slide, Inches(4.5), Inches(4.9), Inches(4.3), Inches(1.2), 
                      RGBColor(0x15, 0x3E, 0x60))
        
        self.add_text_box(slide, Inches(4.5), Inches(5), Inches(4.3), Inches(0.4),
                         '项目状态', font_size=16, bold=True, 
                         color=self.COLOR_WHITE, alignment=PP_ALIGN.CENTER)
        
        self.add_text_box(slide, Inches(4.5), Inches(5.4), Inches(4.3), Inches(0.6),
                         '✅ 开发完成  ✅ 测试通过  🚀 准备上线',
                         font_size=16, color=RGBColor(0x90, 0xCA, 0xF9), 
                         alignment=PP_ALIGN.CENTER)
        
        # 底部
        self.add_text_box(slide, Inches(1.5), Inches(6.6), Inches(10.3), Inches(0.5),
                         '血压守护项目组  |  2026年5月', font_size=14, 
                         color=RGBColor(0x64, 0xB5, 0xF6), alignment=PP_ALIGN.CENTER)
    
    # ==================== 生成PPT ====================
    def generate(self, output_path):
        """生成完整PPT"""
        print('正在生成PPT...')
        
        # 生成所有页面
        self.slide_1_cover()
        print('  ✓ 第1页：封面')
        
        self.slide_2_background()
        print('  ✓ 第2页：社会背景')
        
        self.slide_3_solution()
        print('  ✓ 第3页：解决方案')
        
        self.slide_4_medication()
        print('  ✓ 第4页：服药管理')
        
        self.slide_5_bp_management()
        print('  ✓ 第5页：血压管理')
        
        self.slide_6_alert_family()
        print('  ✓ 第6页：预警与家庭')
        
        self.slide_7_tech()
        print('  ✓ 第7页：技术架构')
        
        self.slide_8_social_value()
        print('  ✓ 第8页：社会价值')
        
        self.slide_9_roadmap()
        print('  ✓ 第9页：实施计划')
        
        self.slide_10_end()
        print('  ✓ 第10页：结束页')
        
        # 保存
        self.prs.save(output_path)
        print(f'\n✅ PPT已保存: {output_path}')


def main():
    """主函数"""
    output_dir = r'E:\开发需要\blood-pressure-guardian'
    output_path = os.path.join(output_dir, '血压守护-项目汇报.pptx')
    
    ppt = BPGuardianPPT()
    ppt.generate(output_path)


if __name__ == '__main__':
    main()
